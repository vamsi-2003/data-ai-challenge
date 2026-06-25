from pptx import Presentation
import os

pptx_path = r"C:\Users\22071\Downloads\Idea Submission Template _ Redrob.pptx"
if not os.path.exists(pptx_path):
    print("PPTX template file not found at:", pptx_path)
else:
    try:
        prs = Presentation(pptx_path)
        print(f"Presentation loaded. Number of slides: {len(prs.slides)}")
        for idx, slide in enumerate(prs.slides):
            print(f"\n--- Slide {idx + 1} ---")
            for shape_idx, shape in enumerate(slide.shapes):
                has_text = shape.has_text_frame
                text = shape.text_frame.text.strip() if has_text else ""
                print(f"Shape {shape_idx}: Name='{shape.name}', HasText={has_text}")
                if text:
                    # Print first 100 characters of text
                    print(f"  Text: {text[:200]}")
    except Exception as e:
        print("Error inspecting PPTX:", e)
