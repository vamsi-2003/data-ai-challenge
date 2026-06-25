#!/usr/bin/env python3
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    BG_COLOR = RGBColor(15, 23, 42)      # Slate 900 (Dark Slate)
    TITLE_COLOR = RGBColor(56, 189, 248)  # Sky 400 (Cyan)
    TEXT_COLOR = RGBColor(241, 245, 249)  # Slate 100 (Off-white)
    ACCENT_COLOR = RGBColor(232, 121, 249) # Fuchsia 400 (Accent)
    MUTED_COLOR = RGBColor(148, 163, 184)  # Slate 400 (Muted)
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800 (Card background)
    
    def apply_bg(slide):
        rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_COLOR
        rect.line.fill.background()
        return slide
        
    def add_slide_header(slide, title_text):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.font.name = 'Segoe UI'
        
    # =========================================================================
    # Slide 1: Title Slide
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    
    # Large Title Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Intelligent Candidate Discovery & Ranking"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = "A Hybrid Semantic & Structured Recruiting System"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_COLOR
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(10)
    
    # Details Box (Bottom left)
    details_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.0))
    tf_det = details_box.text_frame
    tf_det.word_wrap = True
    
    p_team = tf_det.paragraphs[0]
    p_team.text = "Team Name: vamsi krishna"
    p_team.font.size = Pt(18)
    p_team.font.bold = True
    p_team.font.color.rgb = TEXT_COLOR
    p_team.font.name = 'Segoe UI'
    
    p_lead = tf_det.add_paragraph()
    p_lead.text = "Team Leader Name: vamshi krishna vemula(leader)"
    p_lead.font.size = Pt(16)
    p_lead.font.color.rgb = TEXT_COLOR
    p_lead.font.name = 'Segoe UI'
    p_lead.space_before = Pt(5)
    
    p_prob = tf_det.add_paragraph()
    p_prob.text = "Problem Statement: Build an AI candidate ranking system that accurately identifies and ranks the top 100 fits for the Senior AI Engineer role, filtering out keyword stuffers and honeypots under a 5-minute CPU constraint."
    p_prob.font.size = Pt(14)
    p_prob.font.color.rgb = MUTED_COLOR
    p_prob.font.name = 'Segoe UI'
    p_prob.space_before = Pt(10)
    
    # =========================================================================
    # Slide 2: Solution Overview
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Solution Overview")
    
    # Left Box - Proposed Solution
    box_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Proposed Solution"
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_COLOR
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(15)
    
    points_l = [
        "Hybrid Matching: Integrates TF-IDF dense text similarity with a multi-dimensional structured profile scoring model.",
        "Zero-Honeypot Guarantee: Programmatic detection filters out all fake profiles and date-impossible candidate claims.",
        "Keyword-Stuffing Resistance: Evaluates skills strictly by weighting their specified proficiency and duration (beginner, short-term skills get heavily penalized).",
        "Deterministic Tie-Breaker: Sorts equal scores by candidate_id ascending to guarantee full schema validation compatibility."
    ]
    for pt in points_l:
        p = tf_l.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # Right Box - Differentiation
    box_r = slide.shapes.add_textbox(Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Key Differentiators"
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_COLOR
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(15)
    
    points_r = [
        "Operational Realism: Models availability and engagement signals (notice period, response rates, recent active dates) on top of theoretical skill matching.",
        "Fast & Lightweight: Complete end-to-end ranking of 100,000 candidates executes in under 2 minutes on CPU, needing no expensive GPU or LLM API calls.",
        "Factual, Zero-Hallucination Reasoning: Rule-based template compiler dynamically builds 1-2 sentence reasonings based strictly on verified profile facts (no generative LLM hallucinations)."
    ]
    for pt in points_r:
        p = tf_r.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 3: JD Understanding & Candidate Evaluation
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "JD Understanding & Candidate Evaluation")
    
    # Left Column - Key Requirements
    box_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Key JD Requirements"
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_COLOR
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(15)
    
    points_l = [
        "Experience Band: 5-9 years preferred total experience.",
        "Technical Depth: Production experience with sentence-transformers/embeddings and vector databases (Pinecone, Milvus, Qdrant, FAISS).",
        "Evaluation Frameworks: Rigorous offline evaluation metrics (NDCG, MAP, MRR) and online A/B testing background.",
        "Startup Fit: Product/shipper engineering mindset over consulting-only experience (TCS/Wipro/Infosys only backgrounds are filtered)."
    ]
    for pt in points_l:
        p = tf_l.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # Right Column - Signal Importance & Evaluation
    box_r = slide.shapes.add_textbox(Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Signal Importance & Evaluation"
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_COLOR
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(15)
    
    points_r = [
        "Experience Years: Scored highest at 5-9 years (1.0), with linear decay below 5 and slow decay above 9 to favor seniority over juniors.",
        "Role Relevance: Headline and titles scanned for Core AI/ML and Software/Backend keywords. Marketing, HR, or non-technical roles are disqualified.",
        "Behavioral Signals: Multipliers for notice period (<=30 days preferred), activity recency (exponential penalty if inactive >90 days), and responsiveness (response rate & time)."
    ]
    for pt in points_r:
        p = tf_r.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 4: Ranking Methodology
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Ranking Methodology")
    
    # 3-column layout
    col_w = Inches(3.6)
    gap = Inches(0.5)
    
    # Column 1: Retrieval
    box1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), col_w, Inches(4.5))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "1. Retrieval (TF-IDF)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Segoe UI'
    p.space_after = Pt(10)
    
    pts1 = [
        "Combined textual profiles are created from candidate headlines, summaries, job titles, role descriptions, and skills.",
        "TF-IDF vectorization is applied to all candidates.",
        "Cosine similarity computes initial semantic alignment with the JD text."
    ]
    for pt in pts1:
        p = tf1.add_paragraph()
        p.text = "● " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(8)
        
    # Column 2: Structured Scorer
    box2 = slide.shapes.add_textbox(Inches(0.75) + col_w + gap, Inches(1.8), col_w, Inches(4.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "2. Structured Score"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Segoe UI'
    p.space_after = Pt(10)
    
    pts2 = [
        "Experience: Perfect score for 5-9 years band, decaying outside.",
        "Titles: Current/past titles weighted for core AI/ML vs backend SE.",
        "Skills: Evaluated strictly by proficiency (expert=1.0, beginner=0.2) and months of duration to filter keyword stuffers."
    ]
    for pt in pts2:
        p = tf2.add_paragraph()
        p.text = "● " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(8)
        
    # Column 3: Adjustment Multipliers
    box3 = slide.shapes.add_textbox(Inches(0.75) + (col_w + gap)*2, Inches(1.8), col_w, Inches(4.5))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "3. Adjustments"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Segoe UI'
    p.space_after = Pt(10)
    
    pts3 = [
        "Location Multiplier: Noida/Pune preferred (1.0), India Tier-1 relocations (0.95), India other (0.7), international (0.1).",
        "Notice Period Multiplier: <=30 days (1.0), <=90 days (0.8), >90 days (0.5).",
        "Behavioral Multiplier: Multiplies activity recency, response rate/time, and interview completion rates."
    ]
    for pt in pts3:
        p = tf3.add_paragraph()
        p.text = "● " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(8)
        
    # =========================================================================
    # Slide 5: Explainability & Data Validation
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Explainability & Data Validation")
    
    # Left Column - Explainability
    box_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Explainability"
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_COLOR
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(15)
    
    points_l = [
        "Template Compiler: Generates 1-2 sentence rationales referencing candidate's exact profile statistics (e.g. years, location, skills).",
        "Factual Guarantee: Reasoning is built strictly from candidate's JSON profile fields, ensuring zero hallucinated skills or roles.",
        "Variability: Utilizes 4 distinct sentence structures mapped to ranks, ensuring the justifications look custom and natural.",
        "Acknowledge Concerns: Explicitly highlights potential gaps (e.g. high notice periods or experience years slightly outside the band)."
    ]
    for pt in points_l:
        p = tf_l.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # Right Column - Data Validation (Honeypot Filter)
    box_r = slide.shapes.add_textbox(Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Suspicious Profile Filtering (Honeypots)"
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_COLOR
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(15)
    
    points_r = [
        "Flagged Start Date Anomalies: Excludes candidates claiming experience at startups Krutrim or Sarvam AI prior to their actual 2023 founding (e.g. claiming 8 years of experience).",
        "Zero-Duration Skill Flags: Filters out profiles listing expert or advanced skills with duration_months equal to exactly 0.",
        "Consulting-Only Filters: Identifies and penalizes candidates whose entire work histories lie solely in large services companies (TCS, Wipro, Infosys, Cognizant, etc.).",
        "Unrelated Current Titles: Disqualifies keyword stuffers whose current title is unrelated (e.g. Marketing Manager, Accountant)."
    ]
    for pt in points_r:
        p = tf_r.add_paragraph()
        p.text = "●  " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 6: End-to-End Workflow
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "End-to-End Workflow")
    
    # Workflow box
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    steps = [
        "1. Load Data: Read candidates.jsonl line-by-line using a streaming JSON parser, and parse the job description.",
        "2. Parse Activity Date: Detect the maximum activity date in the dataset (2026-05-27) to calculate inactivity recency.",
        "3. Fit TF-IDF Vectorizer: Fit TfidfVectorizer on all candidate text descriptions and compute similarity with JD text.",
        "4. Exclude Honeypots & Disqualified: Run checks on dates, zero-duration expert skills, consulting-only, and unrelated titles.",
        "5. Compute Structured Score: Calculate experience years compatibility, title matches, and Core/Nice-to-have skills alignment.",
        "6. Apply Multipliers: Calculate location, notice period, and behavioral activity multipliers.",
        "7. Sort & Select: Combine TF-IDF and structured scores, sort descending (score) and ascending (candidate_id) for tie-breaks, select top 100.",
        "8. Generate Reasonings & Write CSV: Compile reasoning strings from templates using candidate metrics, and write output to vamsi_krishna.csv."
    ]
    for step in steps:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 7: System Architecture
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "System Architecture")
    
    # Text box for structured description of architecture
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Recruiting Ranking Pipeline"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Segoe UI'
    p.space_after = Pt(15)
    
    arch_details = [
        "Input Data Layer: candidate_profiles (JSONL) + job_description (Markdown)",
        "Pre-filtering Engine (Honeypot Filter -> Consulting Filter -> Non-tech Title Filter)",
        "TF-IDF Vectorization Engine: Fit TfidfVectorizer (Max 15k features) -> Compute Cosine Similarity Matrix",
        "Structured Scorer: Years of Experience (5-9 peak) + Current/Past Title Match + Skill Duration/Proficiency Vector Match",
        "Multiplier Adjuster: Proximity (Noida/Pune) * Notice Period (<=30d) * Behavioral Signals (Activity, Response Rate, Interview Completion)",
        "Combined Scoring Layer: final_score = (tfidf_sim + 0.05) * structured_score * location_mult * notice_mult * behavior_mult",
        "Sorting & Formatting: Sort by (-score, candidate_id) -> Programmatic Reasonings -> Write CSV"
    ]
    for detail in arch_details:
        p = tf.add_paragraph()
        p.text = "➔  " + detail
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(12)
        
    # =========================================================================
    # Slide 8: Results & Performance
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Results & Performance")
    
    # Left Box - Benchmarks
    box_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Execution Benchmarks"
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_COLOR
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(15)
    
    benchmarks = [
        "Total Execution Time: 117.07 seconds (well within the 5-minute limit).",
        "Memory Consumption: Under 500 MB RAM (well within the 16 GB limit).",
        "Honeypot Rate: 0% in Top 100 (Stage 3 filter requires < 10%).",
        "Format Validation: 100% Validated via validate_submission.py."
    ]
    for b in benchmarks:
        p = tf_l.add_paragraph()
        p.text = "● " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(12)
        
    # Right Box - Top 5 Matches Samples
    box_r = slide.shapes.add_textbox(Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Top Rank Highlights"
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_COLOR
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(15)
    
    highlights = [
        "Rank 1 (CAND_0081846): 6.7 years of experience, current title matches Senior AI Engineer, expertise in Elasticsearch & Vector Search, Noida/Pune compatible.",
        "Rank 2 (CAND_0018499): 7.2 years, expert in Weaviate & Pinecone, 15-day notice period.",
        "Rank 3 (CAND_0039754): 16.2 years, expert in Fine-tuning LLMs, Qdrant, OpenSearch, noted experience seniority concern.",
        "Rank 5 (CAND_0099806): 4.6 years, expert in RAG & FAISS, noted experience slightly below range."
    ]
    for h in highlights:
        p = tf_r.add_paragraph()
        p.text = "● " + h
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 9: Technologies Used
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Technologies Used")
    
    # Left Column
    box_l = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Core Stack & Packages"
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_COLOR
    p_lh.font.name = 'Segoe UI'
    p_lh.space_after = Pt(15)
    
    tech_l = [
        "Python: Core language for scripting, feature engineering, and data piping.",
        "Scikit-learn: Utilized TfidfVectorizer for text representation and cosine_similarity for semantic similarity matching.",
        "NumPy & Pandas: Used for vectorized matrix sorting, score calculations, and candidate processing.",
        "PyYAML & PyPDF: Used for managing metadata files and parsing the challenge slide template."
    ]
    for t in tech_l:
        p = tf_l.add_paragraph()
        p.text = "●  " + t
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # Right Column
    box_r = slide.shapes.add_textbox(Inches(6.98), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Design Rationale"
    p_rh.font.size = Pt(22)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_COLOR
    p_rh.font.name = 'Segoe UI'
    p_rh.space_after = Pt(15)
    
    tech_r = [
        "Zero-GPU requirement: Opted for TF-IDF + Cosine similarity over Transformer embeddings to run local scoring within the tight CPU timeline.",
        "Streaming/Iterative parsing: Reads JSONL line-by-line to stay well within the 16 GB memory limit.",
        "Rule-based compilation: Uses deterministic template logic for generating explanations, guaranteeing speed, zero API cost, and zero hallucinations."
    ]
    for t in tech_r:
        p = tf_r.add_paragraph()
        p.text = "●  " + t
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(10)
        
    # =========================================================================
    # Slide 10: Submission Assets
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    add_slide_header(slide, "Submission Assets")
    
    # Asset Details
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The following assets are complete and available in the workspace:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = 'Segoe UI'
    p.space_after = Pt(20)
    
    assets = [
        "GitHub Code Repository: https://github.com/vamsi-2003/data-ai-challenge",
        "HuggingFace Sandbox / Demo Link: https://huggingface.co/spaces/vamsi-2003/data-ai-challenge",
        "Ranked Output File (CSV): vamsi_krishna.csv (contains the top 100 candidates with custom rationales)",
        "Submission Metadata: submission_metadata.yaml (containing team name, leader name, and methodology)",
        "Approach Presentation: vamsi_krishna_approach.pptx (this generated slide deck, ready for PDF export)"
    ]
    for asset in assets:
        p = tf.add_paragraph()
        p.text = "✔  " + asset
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = 'Segoe UI'
        p.space_after = Pt(15)
        
    # =========================================================================
    # Slide 11: Thank You
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    
    # Large Centered Thank You Text
    ty_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "THANK YOU"
    p_ty.font.size = Pt(54)
    p_ty.font.bold = True
    p_ty.font.color.rgb = TITLE_COLOR
    p_ty.font.name = 'Segoe UI'
    p_ty.alignment = 1  # Center alignment
    
    p_sub = tf_ty.add_paragraph()
    p_sub.text = "Build what next India runs on"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = ACCENT_COLOR
    p_sub.font.name = 'Segoe UI'
    p_sub.space_before = Pt(15)
    p_sub.alignment = 1  # Center alignment

    p_cred = tf_ty.add_paragraph()
    p_cred.text = "Team: vamsi krishna  |  Leader: vamshi krishna vemula(leader)"
    p_cred.font.size = Pt(14)
    p_cred.font.color.rgb = MUTED_COLOR
    p_cred.font.name = 'Segoe UI'
    p_cred.space_before = Pt(30)
    p_cred.alignment = 1  # Center alignment
        
    prs.save('vamsi_krishna_approach.pptx')
    print("Successfully created widescreen vamsi_krishna_approach.pptx!")

if __name__ == '__main__':
    create_presentation()
